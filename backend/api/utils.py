import os
import subprocess
import tempfile
import uuid
import sys
from pathlib import Path
from django.conf import settings
from django.core.files import File
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pdf2docx import Converter
import fitz  # PyMuPDF
from pypdf import PdfMerger
from docx import Document
import shutil
import logging
from pptx import Presentation
import openpyxl
import time

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def get_temp_dir():
    """Get the temporary directory path"""
    temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp')
    os.makedirs(temp_dir, exist_ok=True)
    return temp_dir


def get_file_extension(file_path):
    """Get the file extension from a path"""
    return os.path.splitext(file_path)[1][1:].lower()


def convert_to_pdf(input_path, output_path=None):
    """
    Convert various file formats to PDF
    
    Args:
        input_path (str): Path to the input file
        output_path (str, optional): Path for the output PDF file
    
    Returns:
        str: Path to the generated PDF file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.pdf")
    
    file_ext = get_file_extension(input_path)
    
    # Office documents (docx, pptx, xlsx)
    if file_ext in ['docx', 'pptx', 'xlsx']:
        # First, try to use LibreOffice for conversion
        libreoffice_success = False
        try:
            # Check if LibreOffice is available
            user_provided_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
            libreoffice_paths = [
                user_provided_path,
                'soffice',  # Default PATH
                'C:\\Program Files\\LibreOffice\\program\\soffice.exe',  # Common Windows path
                'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',  # 32-bit on 64-bit Windows
                '/usr/bin/soffice',  # Common Linux path
                '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # Mac path
            ]
            
            soffice_path = None
            for path in libreoffice_paths:
                try:
                    # For absolute paths, just check if they exist
                    if os.path.isabs(path):
                        if os.path.exists(path):
                            soffice_path = path
                            break
                    # For 'soffice' in PATH, use 'where' on Windows or 'which' on Unix
                    else:
                        if os.name == 'nt':  # Windows
                            check_cmd = ['where', path]
                        else:  # Unix-like
                            check_cmd = ['which', path]
                        
                        result = subprocess.run(check_cmd, capture_output=True, text=True, check=False)
                        if result.returncode == 0 and result.stdout.strip():
                            soffice_path = result.stdout.strip().splitlines()[0] # take the first one
                            break
                except Exception as e:
                    logger.warning(f"Error checking LibreOffice path {path}: {str(e)}")
            
            if soffice_path:
                logger.info(f"Found LibreOffice at: {soffice_path}")
                # Use LibreOffice for conversion
                cmd = [
                    soffice_path, 
                    '--headless', 
                    '--convert-to', 
                    'pdf',
                    '--outdir', 
                    os.path.dirname(output_path), 
                    input_path
                ]
                
                logger.info(f"Running LibreOffice conversion command: {' '.join(cmd)}")
                result = subprocess.run(cmd, capture_output=True, text=True)
                
                if result.returncode == 0:
                    # LibreOffice keeps the original filename but changes extension
                    original_filename = os.path.basename(input_path)
                    original_name_without_ext = os.path.splitext(original_filename)[0]
                    generated_pdf = os.path.join(os.path.dirname(output_path), f"{original_name_without_ext}.pdf")
                    
                    # Rename to the desired output path
                    if generated_pdf != output_path and os.path.exists(generated_pdf):
                        os.rename(generated_pdf, output_path)
                        libreoffice_success = True
                    elif os.path.exists(generated_pdf):
                        shutil.copyfile(generated_pdf, output_path)
                        libreoffice_success = True
                    else:
                        logger.error(f"LibreOffice conversion failed. Expected output file not found: {generated_pdf}")
                        logger.error(f"Command output: {result.stdout}")
                        logger.error(f"Command error: {result.stderr}")
                else:
                    logger.error(f"LibreOffice conversion failed with return code {result.returncode}")
                    logger.error(f"Command output: {result.stdout}")
                    logger.error(f"Command error: {result.stderr}")
            else:
                logger.warning("LibreOffice not found. Please install it and ensure 'soffice' is in your system's PATH.")
        
        except Exception as e:
            logger.error(f"Error during LibreOffice conversion: {str(e)}")
        
        # If LibreOffice conversion failed, try alternative methods or raise error
        if not libreoffice_success:
            if file_ext == 'docx':
                try:
                    # Try docx2pdf for DOCX files
                    from docx2pdf import convert
                    logger.info(f"Using docx2pdf for conversion of {input_path}")
                    convert(input_path, output_path)
                except Exception as docx_error:
                    logger.error(f"docx2pdf conversion failed: {str(docx_error)}")
                    raise Exception(f"Failed to convert DOCX to PDF. Please ensure LibreOffice is installed correctly for full support.")
            else: # for 'pptx' and 'xlsx'
                raise Exception(f"Failed to convert {file_ext.upper()} to PDF. Please ensure LibreOffice is installed correctly and accessible in the system's PATH.")
    
    # Images (png, jpg, jpeg)
    elif file_ext in ['png', 'jpg', 'jpeg']:
        try:
            img = Image.open(input_path)
            
            # Convert to RGB if needed
            if img.mode != 'RGB':
                img = img.convert('RGB')
                
            # Save as PDF
            img.save(output_path, 'PDF', resolution=100.0)
        except Exception as img_error:
            logger.error(f"Image conversion failed: {str(img_error)}")
            raise Exception(f"Failed to convert image to PDF: {str(img_error)}")
    
    # Text files
    elif file_ext == 'txt':
        try:
            # Create a PDF from text
            c = canvas.Canvas(output_path, pagesize=letter)
            text_file = open(input_path, 'r', encoding='utf-8', errors='ignore')
            
            y = 750  # Starting y position
            line_height = 14
            
            for line in text_file:
                if y < 50:  # Create a new page when reaching bottom
                    c.showPage()
                    y = 750
                
                c.drawString(50, y, line.strip())
                y -= line_height
            
            text_file.close()
            c.save()
        except Exception as txt_error:
            logger.error(f"Text file conversion failed: {str(txt_error)}")
            raise Exception(f"Failed to convert TXT to PDF: {str(txt_error)}")
    
    # PDF files (just copy)
    elif file_ext == 'pdf':
        try:
            shutil.copyfile(input_path, output_path)
        except Exception as copy_error:
            logger.error(f"PDF copy failed: {str(copy_error)}")
            raise Exception(f"Failed to copy PDF file: {str(copy_error)}")
    
    else:
        raise ValueError(f"Unsupported file format for conversion: {file_ext}")
    
    return output_path


def pdf_to_docx(input_path, output_path=None):
    """
    Convert PDF to DOCX
    
    Args:
        input_path (str): Path to the input PDF file
        output_path (str, optional): Path for the output DOCX file
    
    Returns:
        str: Path to the generated DOCX file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.docx")
    
    # Check if input is actually a PDF
    if get_file_extension(input_path) != 'pdf':
        raise ValueError("Input file must be a PDF")
    
    # Convert PDF to DOCX
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()
    
    return output_path


def pdf_to_txt(input_path, output_path=None):
    """
    Convert PDF to TXT
    
    Args:
        input_path (str): Path to the input PDF file
        output_path (str, optional): Path for the output TXT file
    
    Returns:
        str: Path to the generated TXT file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.txt")
    
    # Check if input is actually a PDF
    if get_file_extension(input_path) != 'pdf':
        raise ValueError("Input file must be a PDF")
    
    # Extract text from PDF
    pdf_document = fitz.open(input_path)
    with open(output_path, 'w', encoding='utf-8') as txt_file:
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            txt_file.write(page.get_text())
            txt_file.write('\n\n--- Page Break ---\n\n')
    
    pdf_document.close()
    
    return output_path


def pdf_to_pptx(input_path, output_path=None):
    """
    Convert PDF to PPTX. Each PDF page becomes an image on a slide.
    
    Args:
        input_path (str): Path to the input PDF file.
        output_path (str, optional): Path for the output PPTX file.
        
    Returns:
        str: Path to the generated PPTX file.
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.pptx")

    if get_file_extension(input_path) != 'pdf':
        raise ValueError("Input file must be a PDF")

    pdf_document = fitz.open(input_path)
    prs = Presentation()
    
    # Use a temp directory for images
    temp_dir = tempfile.mkdtemp()

    try:
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            pix = page.get_pixmap(dpi=150) # Use a reasonable DPI
            img_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(img_path)

            blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Set slide dimensions from PDF page dimensions
            prs.slide_width = int(page.rect.width * 914400 / 72) # Convert points to EMU
            prs.slide_height = int(page.rect.height * 914400 / 72)

            # Add image to fill the slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    finally:
        pdf_document.close()
        # Clean up temporary directory
        shutil.rmtree(temp_dir, ignore_errors=True)

    prs.save(output_path)
    return output_path


def merge_pdf_files(file_paths, output_path=None):
    """
    Merge multiple PDF files into one
    
    Args:
        file_paths (list): List of paths to PDF files to merge
        output_path (str, optional): Path for the output merged PDF
    
    Returns:
        str: Path to the merged PDF file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.pdf")
    
    logger.info(f"Merging {len(file_paths)} PDF files to {output_path}")
    
    # Check if all files are PDFs and exist
    for file_path in file_paths:
        if not os.path.exists(file_path):
            raise ValueError(f"PDF file not found: {file_path}")
        if get_file_extension(file_path) != 'pdf':
            raise ValueError(f"File '{file_path}' is not a PDF")
    
    try:
        # Merge PDFs
        merger = PdfMerger()
        for i, file_path in enumerate(file_paths):
            try:
                merger.append(file_path)
                logger.info(f"Added PDF {i+1}/{len(file_paths)}: {file_path}")
            except Exception as e:
                raise ValueError(f"Error adding PDF file {file_path}: {str(e)}")
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Write merged PDF to output file
        merger.write(output_path)
        merger.close()
        
        if not os.path.exists(output_path):
            raise ValueError(f"Failed to create merged PDF at {output_path}")
            
        logger.info(f"Successfully merged PDFs to {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Error merging PDF files: {str(e)}")
        # Clean up partial output file if it exists
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except:
                pass
        raise Exception(f"Failed to merge PDF files: {str(e)}")


def merge_docx_files(file_paths, output_path=None):
    """
    Merge multiple DOCX files into one
    
    Args:
        file_paths (list): List of paths to DOCX files to merge
        output_path (str, optional): Path for the output merged DOCX
    
    Returns:
        str: Path to the merged DOCX file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.docx")
    
    logger.info(f"Merging {len(file_paths)} DOCX files to {output_path}")
    
    # Check if all files are DOCX and exist
    for file_path in file_paths:
        if not os.path.exists(file_path):
            raise ValueError(f"DOCX file not found: {file_path}")
        if get_file_extension(file_path) != 'docx':
            raise ValueError(f"File '{file_path}' is not a DOCX")
    
    try:
        # Create a new document for the merged content
        merged_doc = Document()
        
        # Process each document
        for i, file_path in enumerate(file_paths):
            try:
                logger.info(f"Processing DOCX {i+1}/{len(file_paths)}: {file_path}")
                doc = Document(file_path)
                
                # Add a page break between documents (except for the first one)
                if i > 0:
                    merged_doc.add_page_break()
                
                # Copy content from each document
                for element in doc.element.body:
                    merged_doc.element.body.append(element)
                    
            except Exception as e:
                raise ValueError(f"Error processing DOCX file {file_path}: {str(e)}")
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Save the merged document
        merged_doc.save(output_path)
        
        if not os.path.exists(output_path):
            raise ValueError(f"Failed to create merged DOCX at {output_path}")
            
        logger.info(f"Successfully merged DOCX files to {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Error merging DOCX files: {str(e)}")
        # Clean up partial output file if it exists
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except:
                pass
        raise Exception(f"Failed to merge DOCX files: {str(e)}")


def merge_pptx_files(file_paths, output_path=None):
    """
    Merge multiple PPTX files into one
    
    Args:
        file_paths (list): List of paths to PPTX files to merge
        output_path (str, optional): Path for the output merged PPTX
    
    Returns:
        str: Path to the merged PPTX file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.pptx")
    
    # Check if all files are PPTX
    for file_path in file_paths:
        if not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")
        if get_file_extension(file_path) != 'pptx':
            raise ValueError("All input files must be PPTX files")
    
    # For PPTX, we'll use Python-pptx library
    # Create a new presentation
    merged_pres = Presentation()
    
    # Remove default slide
    if len(merged_pres.slides) > 0:
        xml_slides = merged_pres.slides._sldIdLst
        slides = list(xml_slides)
        for slide in slides:
            xml_slides.remove(slide)
    
    # Loop through each presentation to add its slides
    for file_path in file_paths:
        try:
            # Open the presentation
            pres = Presentation(file_path)
            
            # Loop through all slides and add them to merged presentation
            for slide in pres.slides:
                # Copy slide and its layout
                slide_layout = merged_pres.slide_layouts[0]  # Default to first layout
                new_slide = merged_pres.slides.add_slide(slide_layout)
                
                # Copy content by element XML
                for shape in slide.shapes:
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                    
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {str(e)}")
            raise Exception(f"Failed to process PPTX file: {str(e)}")
    
    # Save the merged presentation
    try:
        merged_pres.save(output_path)
    except Exception as e:
        logger.error(f"Error saving merged PPTX: {str(e)}")
        raise Exception(f"Failed to save merged PPTX: {str(e)}")
    
    return output_path


def merge_files(file_paths, output_filename, file_type):
    """
    Merge files of the same type
    
    Args:
        file_paths (list): List of file paths to merge
        output_filename (str): Name for the output file
        file_type (str): Type of files being merged (pdf, docx, pptx)
    
    Returns:
        str: Path to the merged file
    """
    logger.info(f"Starting merge operation for {len(file_paths)} files of type {file_type}")
    
    # Basic validation
    if not file_paths:
        raise ValueError("No files provided for merging")
    
    if len(file_paths) < 2:
        raise ValueError("At least two files are required for merging")
    
    # Ensure all files exist and are of the correct type
    for file_path in file_paths:
        if not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")
        
        if get_file_extension(file_path).lower() != file_type.lower():
            raise ValueError(f"File '{file_path}' is not a {file_type} file")
    
    # Sanitize the output filename (remove any directory traversal attempts)
    safe_filename = os.path.basename(output_filename)
    if safe_filename != output_filename:
        logger.warning(f"Output filename sanitized from {output_filename} to {safe_filename}")
    
    # Determine output path
    output_path = os.path.join(get_temp_dir(), f"{safe_filename}.{file_type}")
    logger.info(f"Output path for merged file: {output_path}")
    
    try:
        # Merge based on file type
        if file_type.lower() == 'pdf':
            return merge_pdf_files(file_paths, output_path)
        elif file_type.lower() == 'docx':
            return merge_docx_files(file_paths, output_path)
        elif file_type.lower() == 'pptx':
            return merge_pptx_files(file_paths, output_path)
        else:
            raise ValueError(f"Unsupported file type for merging: {file_type}")
    except Exception as e:
        logger.error(f"Error during merge operation: {str(e)}")
        # Make sure we don't leave a partial output file
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except:
                pass
        raise


def clean_temp_files():
    """Clean up temporary files older than 1 hour"""
    import time
    temp_dir = get_temp_dir()
    current_time = time.time()
    
    for filename in os.listdir(temp_dir):
        file_path = os.path.join(temp_dir, filename)
        
        # If file is older than 1 hour (3600 seconds)
        if os.path.isfile(file_path) and current_time - os.path.getmtime(file_path) > 3600:
            try:
                os.remove(file_path)
            except:
                pass


def merge_images_to_pdf(image_paths, output_path=None):
    """
    Merge multiple images into a single PDF file, preserving the order
    
    Args:
        image_paths (list): List of paths to image files
        output_path (str, optional): Path for the output PDF file
    
    Returns:
        str: Path to the generated PDF file
    """
    if output_path is None:
        output_path = os.path.join(get_temp_dir(), f"{uuid.uuid4()}.pdf")
    
    # Check if all files are images
    for image_path in image_paths:
        ext = get_file_extension(image_path)
        if ext not in ['png', 'jpg', 'jpeg']:
            raise ValueError(f"Unsupported file format: {ext}. Only PNG, JPG, and JPEG are supported")
    
    # If there's only one image, convert it directly to PDF
    if len(image_paths) == 1:
        img = Image.open(image_paths[0])
        
        # Convert to RGB if needed
        if img.mode != 'RGB':
            img = img.convert('RGB')
            
        # Save as PDF
        img.save(output_path, 'PDF', resolution=100.0)
        return output_path
    
    # For multiple images, create a PDF with multiple pages
    images = []
    for image_path in image_paths:
        img = Image.open(image_path)
        # Convert to RGB if needed
        if img.mode != 'RGB':
            img = img.convert('RGB')
        images.append(img)
    
    # Save the first image as PDF
    images[0].save(
        output_path, 
        'PDF', 
        resolution=100.0,
        save_all=True, 
        append_images=images[1:]
    )
    
    return output_path

# Function to process files without database dependency
def process_file_without_db(uploaded_file, operation):
    """
    Process a file without requiring database access
    
    Args:
        uploaded_file: The uploaded file object
        operation: The operation to perform (e.g., 'convert_to_pdf')
    
    Returns:
        tuple: (output_path, output_filename)
    """
    try:
        # Create temp directory if it doesn't exist
        temp_dir = get_temp_dir()
        
        # Save uploaded file to temp location
        timestamp = int(time.time())
        temp_input_path = os.path.join(temp_dir, f"input_{timestamp}_{uploaded_file.name}")
        
        # Ensure the directory exists
        os.makedirs(os.path.dirname(temp_input_path), exist_ok=True)
        
        # Save the uploaded file
        with open(temp_input_path, 'wb+') as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)
        
        # Get file extension
        file_name = uploaded_file.name
        extension = file_name.split('.')[-1].lower()
        
        # Process based on operation
        if operation == 'convert_to_pdf':
            # Check if file is already a PDF
            if extension == 'pdf':
                return temp_input_path, file_name
            
            # Convert to PDF
            output_path = convert_to_pdf(temp_input_path)
            output_filename = f"{os.path.splitext(file_name)[0]}.pdf"
            
        elif operation == 'pdf_to_docx':
            # Check if file is a PDF
            if extension != 'pdf':
                raise ValueError('Only PDF files can be converted to DOCX')
            
            # Convert PDF to DOCX
            output_path = pdf_to_docx(temp_input_path)
            output_filename = f"{os.path.splitext(file_name)[0]}.docx"
            
        elif operation == 'pdf_to_txt':
            # Check if file is a PDF
            if extension != 'pdf':
                raise ValueError('Only PDF files can be converted to TXT')
            
            # Convert PDF to TXT
            output_path = pdf_to_txt(temp_input_path)
            output_filename = f"{os.path.splitext(file_name)[0]}.txt"
            
        elif operation == 'pdf_to_pptx':
            if extension != 'pdf':
                raise ValueError('Only PDF files can be converted to PPTX')
            
            # Convert PDF to PPTX
            output_path = pdf_to_pptx(temp_input_path)
            output_filename = f"{os.path.splitext(file_name)[0]}.pptx"
            
        else:
            raise ValueError(f'Unsupported operation: {operation}')
        
        return output_path, output_filename
        
    except Exception as e:
        logger.error(f"Error processing file without DB: {str(e)}")
        raise e
    finally:
        # Clean up input file
        try:
            if 'temp_input_path' in locals() and os.path.exists(temp_input_path):
                os.remove(temp_input_path)
        except Exception as e:
            logger.error(f"Error cleaning up temp file: {str(e)}")

# Function to process multiple images to PDF without database dependency
def process_images_to_pdf_without_db(files, output_filename):
    """
    Process multiple images to create a PDF without requiring database access
    
    Args:
        files: List of uploaded file objects
        output_filename: Name for the output file
    
    Returns:
        tuple: (output_path, output_filename)
    """
    try:
        # Create temp directory if it doesn't exist
        temp_dir = get_temp_dir()
        timestamp = int(time.time())
        
        # Save uploaded files to temp location
        file_paths = []
        for i, file in enumerate(files):
            temp_path = os.path.join(temp_dir, f"img_{timestamp}_{i}_{file.name}")
            
            # Ensure the directory exists
            os.makedirs(os.path.dirname(temp_path), exist_ok=True)
            
            # Save the uploaded file
            with open(temp_path, 'wb+') as destination:
                for chunk in file.chunks():
                    destination.write(chunk)
            
            file_paths.append(temp_path)
        
        # Merge images to PDF
        output_path = merge_images_to_pdf(file_paths)
        final_output_filename = f"{output_filename}.pdf"
        
        return output_path, final_output_filename
        
    except Exception as e:
        logger.error(f"Error processing images to PDF without DB: {str(e)}")
        raise e
    finally:
        # Clean up input files
        for path in file_paths if 'file_paths' in locals() else []:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception as e:
                logger.error(f"Error cleaning up temp file {path}: {str(e)}")

# Function to merge files without database dependency
def merge_files_without_db(files, output_filename, file_type):
    """
    Merge multiple files without requiring database access
    
    Args:
        files: List of uploaded file objects
        output_filename: Name for the output file
        file_type: Type of files being merged (pdf, docx, pptx)
    
    Returns:
        tuple: (output_path, output_filename)
    """
    try:
        # Create temp directory if it doesn't exist
        temp_dir = get_temp_dir()
        timestamp = int(time.time())
        
        # Save uploaded files to temp location
        file_paths = []
        for i, file in enumerate(files):
            temp_path = os.path.join(temp_dir, f"merge_{timestamp}_{i}_{file.name}")
            
            # Ensure the directory exists
            os.makedirs(os.path.dirname(temp_path), exist_ok=True)
            
            # Save the uploaded file
            with open(temp_path, 'wb+') as destination:
                for chunk in file.chunks():
                    destination.write(chunk)
            
            file_paths.append(temp_path)
        
        # Merge files
        output_path = merge_files(file_paths, output_filename, file_type)
        final_output_filename = f"{output_filename}.{file_type}"
        
        return output_path, final_output_filename
        
    except Exception as e:
        logger.error(f"Error merging files without DB: {str(e)}")
        raise e
    finally:
        # Clean up input files
        for path in file_paths if 'file_paths' in locals() else []:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception as e:
                logger.error(f"Error cleaning up temp file {path}: {str(e)}") 